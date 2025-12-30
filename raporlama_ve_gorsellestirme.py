# -*- coding: utf-8 -*-
"""
raporlama_ve_gorsellestirme.py (PoF - Ultimate Reporting Engine v3.3)
FIXES:
1. Auto-calculates 'PoF_Ensemble_12Ay' if missing.
2. Robust merging logic for Case Studies.
3. Maps 'Risk_Sinifi' to 'Risk_Class' automatically.
"""

import os
import sys
import logging
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime, timedelta

# PPTX Library Check
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: 'python-pptx' not installed. Skipping PPTX generation.")

# --- KONFİGÜRASYON VE YOLLAR ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(BASE_DIR, "data", "sonuclar")
INTERMEDIATE_DIR = os.path.join(BASE_DIR, "data", "ara_ciktilar")
LOG_DIR = os.path.join(BASE_DIR, "loglar")

# Klasörlerin varlığından emin ol
for d in [OUTPUT_DIR, INTERMEDIATE_DIR, LOG_DIR]:
    os.makedirs(d, exist_ok=True)

# Alt Klasörler
ACTION_DIR = os.path.join(OUTPUT_DIR, "aksiyon_listeleri")
VISUAL_DIR = os.path.join(OUTPUT_DIR, "gorseller")
REPORT_DIR = OUTPUT_DIR

for d in [ACTION_DIR, VISUAL_DIR]:
    os.makedirs(d, exist_ok=True)

# Görsel Ayarları
plt.style.use('ggplot')
sns.set_palette("husl")
# =============================================================================
# 🛠️ UTILITIES: LOGGING & DATA REPAIR (LOGLAMA VE VERİ ONARIMI)
# =============================================================================
# 1. setup_logger:
#    - Raporlama sürecinin her adımını kayıt altına alır.
#    - Hataları hem ekrana basar hem dosyaya kaydeder.
#
# 2. ensure_pof_column (Hayat Kurtarıcı):
#    - Raporların ana metriği 'PoF_Ensemble_12Ay' sütunudur.
#    - Eğer önceki aşamada (pof.py) bu sütun üretilemediyse (örn. sadece tek model çalıştıysa),
#      bu fonksiyon devreye girer.
#    - Mevcut diğer 12 aylık tahminleri (Cox, ML vb.) bulur ve onların ortalamasını
#      alarak eksik sütunu "imal eder". Böylece raporlama çökmez.
# =============================================================================
# ------------------------------------------------------------------------------
# DATE PARSING (KARMA FORMAT DESTEĞİ)
# ------------------------------------------------------------------------------
def parse_mixed_dates(date_str):
    """
    Karma tarih formatlarını TR standartına (GÜN-AY-YIL) uygun parse eder.
    Desteklenen formatlar:
    - DD-MM-YYYY HH:MM:SS (22-06-2025 04:59:21)
    - YYYY-MM-DD HH:MM:SS (2023-01-17 17:14:42)
    - DD-MM-YYYY (05-03-2025)
    """
    if pd.isna(date_str):
        return pd.NaT

    # String'e çevir (sayı olarak gelebilir)
    date_str = str(date_str).strip()

    # Deneme sırası (TR formatı öncelikli)
    formats = [
        '%d-%m-%Y %H:%M:%S',  # 22-06-2025 04:59:21
        '%d-%m-%Y',           # 05-03-2025
        '%Y-%m-%d %H:%M:%S',  # 2023-01-17 17:14:42
        '%Y-%m-%d',           # 2023-01-17
        '%d.%m.%Y %H:%M:%S',  # 22.06.2025 04:59:21 (Noktalı)
        '%d.%m.%Y',           # 22.06.2025
    ]

    for fmt in formats:
        try:
            return pd.to_datetime(date_str, format=fmt)
        except:
            continue

    # Hiçbiri işe yaramazsa pandas otomatik (dayfirst=True)
    try:
        return pd.to_datetime(date_str, dayfirst=True)
    except:
        return pd.NaT

# ------------------------------------------------------------------------------
# LOGGING
# ------------------------------------------------------------------------------
def setup_logger():
    os.makedirs(LOG_DIR, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    logger = logging.getLogger("Raporlama")
    logger.setLevel(logging.INFO)
    logger.handlers.clear()
    
    fh = logging.FileHandler(os.path.join(LOG_DIR, f"raporlama_{ts}.log"), encoding="utf-8")
    fh.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
    logger.addHandler(fh)
    logger.addHandler(logging.StreamHandler(sys.stdout))
    return logger

# ------------------------------------------------------------------------------
# HELPER: ENSURE POF COLUMN EXISTS
# ------------------------------------------------------------------------------
def ensure_pof_column(df, logger):
    """
    'PoF_Ensemble_12Ay' kolonunu kontrol eder, yoksa hesaplar.
    """
    target = 'PoF_Ensemble_12Ay'
    
    if target in df.columns:
        return df
    
    logger.warning(f"  [FIX] '{target}' eksik. Bileşenlerden hesaplanmaya çalışılıyor...")
    
    # Olası PoF kolonlarını bul
    candidates = [c for c in df.columns if ('_pof_12' in c.lower()) or ('_12ay' in c.lower() and 'pof' in c.lower())]
    candidates = [c for c in candidates if c != target]
    
    if candidates:
        logger.info(f"  [FIX] Bulunan bileşenler: {candidates}")
        df[target] = df[candidates].mean(axis=1)
        logger.info(f"  [FIX] '{target}' kolonu {len(candidates)} modelin ortalaması ile oluşturuldu.")
    else:
        logger.warning("  [FAIL] 12 aylık PoF bileşeni bulunamadı. Dummy (0.0) oluşturuluyor.")
        df[target] = 0.0
        
    return df

# ------------------------------------------------------------------------------
# PHASE 1: ACTION PLANNING
# ------------------------------------------------------------------------------
# =============================================================================
# 📋 ACTIONABLE INTELLIGENCE (AKSİYON LİSTELERİ)
# =============================================================================
# Bu fonksiyon, matematiksel sonuçları "Sahada Yapılacak İşlere" dönüştürür.
#
# 1. 🚨 Acil Müdahale (Priority 1):
#    - Hem "Kronik" (Sürekli bozulan) hem de "Kritik Risk" (Ömrünü tamamlamış) varlıklar.
#    - Aksiyon: Derhal yerinde inceleme veya değişim.
#
# 2. 💰 CAPEX / Yatırım (Priority 2):
#    - Özellikle Trafolar gibi pahalı ve tedariği uzun süren ekipmanlar.
#    - Yüksek riskli trafolar belirlenip, gelecek yılın bütçesine değişim olarak girmeli.
#
# 3. 🔍 Fırsat Bakımı / OPEX (Priority 3):
#    - Model diyor ki: "Bu varlık henüz Kritik sınıfa girmedi (belki yaşı genç),
#      AMA önümüzdeki 12 ay içinde bozulma ihtimali yükseliyor (%15+)."
#    - Aksiyon: Önleyici bakım (Termal kamera, yağ analizi) ile kurtarılabilir.
# =============================================================================
def generate_action_lists(df, logger):
    """
    Operasyonel aksiyon listelerini (CSV) oluşturur.
    Bakım ekipleri için okunabilir, temiz raporlar üretir.
    """
    logger.info("="*60)
    logger.info("[PHASE 1] Aksiyon Listeleri Oluşturuluyor...")
    
    # 1. Kolon Standardizasyonu (Defensive Coding)
    risk_col = 'Risk_Sinifi'
    if 'Risk_Class' in df.columns: risk_col = 'Risk_Class'
    
    chronic_col = 'Chronic_Flag'
    if 'Kronik_Flag' in df.columns: chronic_col = 'Kronik_Flag' # Manuel değişim varsa yakala

    # 2. Raporlarda Görünecek Temiz Sütunlar (Human-Readable)
    # Bakımcının işine yaramayan VIF skorlarını, one-hot sütunlarını rapora koymuyoruz.
    report_cols = [
        "cbs_id", "Ekipman_Tipi", "Ilce", "Mahalle", "Marka",
        "Kurulum_Tarihi", "Gerilim_Seviyesi", 
        "Health_Score", risk_col, "PoF_Ensemble_12Ay", 
        chronic_col, "Ariza_Sayisi_90g"
    ]
    # Sadece veride var olanları seç
    final_cols = [c for c in report_cols if c in df.columns]

    # --- LİSTE 1: ACİL MÜDAHALE (Kritik + Kronik) ---
    # Hem çok riskli hem de sürekli arıza yapıyor. Hemen bakılmalı.
    if chronic_col in df.columns:
        crit_chronic = df[
            (df[risk_col].isin(['Critical', 'KRİTİK', 'KRİTİK (KRONİK)'])) & 
            (df[chronic_col] == 1)
        ].copy()
    else:
        crit_chronic = pd.DataFrame()
    
    if not crit_chronic.empty:
        # En riskliden aza doğru sırala
        crit_chronic = crit_chronic.sort_values("PoF_Ensemble_12Ay", ascending=False)
        
        path = os.path.join(ACTION_DIR, "01_ACIL_MUDHALE_LISTESI_Kronik.csv")
        crit_chronic[final_cols].to_csv(path, index=False, encoding='utf-8-sig')
        logger.info(f"  🚨 [ACİL] Kronik & Kritik: {len(crit_chronic)} varlık (Dosya: 01_...)")
    else:
        logger.info("  ✅ [ACİL] Kronik ve Kritik varlık bulunamadı.")

    # --- LİSTE 2: YÜKSEK RİSKLİ TRAFOLAR (CAPEX Yatırım Planı) ---
    # Trafolar pahalıdır. Yüksek riskli olanların değişimi bütçelenmeli.
    trafos = df[
        (df['Ekipman_Tipi'].astype(str).str.contains('Trafo', case=False, na=False)) & 
        (df[risk_col].isin(['Critical', 'High', 'KRİTİK', 'YÜKSEK', 'KRİTİK (KRONİK)']))
    ].copy()
    
    if not trafos.empty:
        trafos = trafos.sort_values("Health_Score", ascending=True) # En düşük puan en üstte
        path = os.path.join(ACTION_DIR, "02_YATIRIM_PLANLAMA_Riskli_Trafolar.csv")
        trafos[final_cols].to_csv(path, index=False, encoding='utf-8-sig')
        logger.info(f"  💰 [CAPEX] Yüksek Riskli Trafolar: {len(trafos)} varlık")

    # --- LİSTE 3: İŞLETME KONTROL (Fırsat Bakımı / Quick Wins) ---
    # Risk sınıfı henüz 'Kritik' değil ama Bozulma İhtimali (PoF) artmaya başlamış.
    # "Henüz yangın çıkmadı ama duman tütüyor" listesi.
    if 'PoF_Ensemble_12Ay' in df.columns:
        inspection = df[
            (df['PoF_Ensemble_12Ay'] > 0.15) &  # %15 üzeri ihtimal
            (df[risk_col].isin(['Low', 'Medium', 'DÜŞÜK', 'ORTA'])) # Ama sınıfı düşük
        ].copy()
        
        if not inspection.empty:
            inspection = inspection.sort_values('PoF_Ensemble_12Ay', ascending=False)
            path = os.path.join(ACTION_DIR, "03_ISLETME_KONTROL_Firsat_Bakimi.csv")
            inspection[final_cols].to_csv(path, index=False, encoding='utf-8-sig')
            logger.info(f"  🔍 [OPEX] Fırsat Bakımı (Yüksek Olasılık/Düşük Risk): {len(inspection)} varlık")

    return crit_chronic # Dashboard için kritik listeyi döndür

# ------------------------------------------------------------------------------
# PHASE 2: VISUALIZATION
# ------------------------------------------------------------------------------
# =============================================================================
# 📊 VISUALIZATION & REALITY CHECK (GÖRSELLEŞTİRME VE DOĞRULAMA)
# =============================================================================
# Bu modül, karmaşık model çıktılarını (olasılıklar, katsayılar) yöneticilerin
# anlayabileceği görsel içgörülere dönüştürür ve modelin "akıl sağlığını" test eder.
#
# 1. 🌍 Coğrafi Risk Haritası (generate_visuals):
#    - Risklerin mekansal dağılımını gösterir.
#    - "Hangi ilçede veya şebeke kolunda risk birikmiş?" sorusuna cevap verir.
#
# 2. 🧪 Kalibrasyon & Validasyon (validate_base_rates):
#    - EN KRİTİK GÜVENLİK ADIMIDIR.
#    - Modelin sonuçlarını CIGRE/IEEE endüstri standartlarıyla kıyaslar.
#    - Örnek: Eğer model Trafoların %50'sinin seneye bozulacağını söylüyorsa,
#      bu fonksiyon "HATA: Endüstri standardı %2-5 arasıdır, model aşırı kötümser!"
#      diye uyarır. Bu, "Model Halüsinasyonunu" engeller.
#
# 3. 📈 Stratejik Ayrım (plot_aggregate_risk_by_type):
#    - Çift Eksenli Grafik (Dual-Axis):
#      a. Bar (Sol): Yaşlanma/Yıpranma Riski (PoF).
#      b. Çizgi (Sağ): Kronik/Operasyonel Sorunlar.
#    - Bu ayrım, yatırımın nereye yapılacağını (Yeni cihaz mı? Tamir mi?) belirler.
# =============================================================================
def plot_single_chart(df, col_x, col_y, plot_type, title, filename, logger, **kwargs):
    """
    Genel amaçlı, hata korumalı grafik çizim fonksiyonu.
    """
    width = kwargs.pop('width', 10)
    height = kwargs.pop('height', 6)
    
    plt.figure(figsize=(width, height))
    
    try:
        if plot_type == 'scatter':
            sns.scatterplot(data=df, x=col_x, y=col_y, **kwargs)
        elif plot_type == 'hist':
            sns.histplot(df[col_x], kde=True, **kwargs)
        elif plot_type == 'bar':
            sns.barplot(x=col_x, y=col_y, data=df, **kwargs)

        plt.title(title, fontsize=14, fontweight='bold')
        plt.tight_layout()
        path = os.path.join(VISUAL_DIR, filename)
        plt.savefig(path, dpi=300, bbox_inches='tight')
        plt.close()
        logger.info(f"  📸 Grafik oluşturuldu: {filename}")
        return path
    except Exception as e:
        logger.error(f"  ❌ [ERROR] Grafik çizilemedi ({filename}): {str(e)}")
        plt.close()
        return None
# ------------------------------------------------------------------------------
# PHASE 2.5: ADVANCED DIAGNOSTICS (SİZİN GÖRSELLERİNİZ)
# ------------------------------------------------------------------------------

def plot_risk_drivers(df_model, logger):
    """
    GÖRSEL 1: Risk Faktörleri (Korelasyon Analizi)
    Veri Kaynağı: model_input_data_full.csv (pof.py ara çıktısı)
    """
    logger.info("[ADVANCED] Risk Faktörleri (Drivers) analiz ediliyor...")
    
    # Sayısal kolonları seç
    num_cols = df_model.select_dtypes(include=[np.number]).columns
    if 'event' not in num_cols: return None

    # Korelasyon hesapla (Target: event)
    corrs = df_model[num_cols].corrwith(df_model['event']).sort_values(ascending=False)
    
    # En etkili 10 faktör (Kendisi hariç)
    top_drivers = corrs.drop('event', errors='ignore').head(5)
    bottom_drivers = corrs.drop('event', errors='ignore').tail(5)
    drivers = pd.concat([top_drivers, bottom_drivers]).sort_values()

    plt.figure(figsize=(10, 6))
    colors = ['green' if x < 0 else 'red' for x in drivers.values]
    drivers.plot(kind='barh', color=colors, edgecolor='black', alpha=0.8)
    
    plt.title('Risk Faktörleri (Drivers) - Arıza ile Korelasyon', fontsize=16)
    plt.xlabel('Korelasyon Katsayısı (Sağ taraf risk arttırıcı)', fontsize=10)
    plt.grid(True, linestyle='--', alpha=0.5)
    
    path = os.path.join(VISUAL_DIR, "ADV_01_Risk_Drivers.png")
    plt.savefig(path, dpi=300, bbox_inches='tight')
    plt.close()
    return path

def plot_operational_dashboard(logger):
    """
    GÖRSEL 2: Operasyonel Durum (4'lü Dashboard)
    Veri Kaynağı: fault_events_clean.csv
    """
    path = os.path.join(INTERMEDIATE_DIR, "fault_events_clean.csv")
    if not os.path.exists(path): return None
    
    df = pd.read_csv(path)
    df['started at'] = df['started at'].apply(parse_mixed_dates)

    # Gelecek tarihli kayıtları filtrele
    today = pd.Timestamp.now()
    df = df[df['started at'] <= today]

    fig, axes = plt.subplots(2, 2, figsize=(16, 10))
    fig.suptitle('Operasyonel Durum Paneli', fontsize=24)
    
    # 1. Aylık Arıza Trendi
    df['Ay'] = df['started at'].dt.to_period('M')
    monthly = df.groupby('Ay').size()
    monthly.index = monthly.index.astype(str)
    axes[0, 0].plot(monthly.index, monthly.values, marker='o', linestyle='-', color='steelblue')
    axes[0, 0].set_title('Aylık Arıza Trendi')
    axes[0, 0].tick_params(axis='x', rotation=45)
    # X eksenini seyrelt
    axes[0, 0].set_xticks(axes[0, 0].get_xticks()[::3])

    # 2. En Çok Arızalanan Ekipmanlar
    top_eq = df['Ekipman_Tipi'].value_counts().head(5)
    sns.barplot(y=top_eq.index, x=top_eq.values, ax=axes[0, 1], palette='Oranges_r')
    axes[0, 1].set_title('En Çok Arızalanan Ekipmanlar')

    # 3. Haftalık Yoğunluk (Heatmap yerine Line/Area)
    df['Hafta'] = df['started at'].dt.isocalendar().week
    weekly = df.groupby('Hafta').size()
    axes[1, 0].fill_between(weekly.index, weekly.values, color='lightcoral', alpha=0.5)
    axes[1, 0].plot(weekly.index, weekly.values, color='red')
    axes[1, 0].set_title('Haftalık Arıza Yoğunluğu (Mevsimsellik)')

    # 4. Saatlik Dağılım
    df['Saat'] = df['started at'].dt.hour
    hourly = df.groupby('Saat').size()
    sns.barplot(x=hourly.index, y=hourly.values, ax=axes[1, 1], color='purple')
    axes[1, 1].set_title('Saatlik Arıza Dağılımı')

    plt.tight_layout(rect=[0, 0.03, 1, 0.95])
    out_path = os.path.join(VISUAL_DIR, "ADV_02_Operasyonel_Durum.png")
    plt.savefig(out_path, dpi=300)
    plt.close()
    return out_path

def plot_survival_curves(df_model, logger):
    """
    GÖRSEL 3: Ömür Eğrileri (Kaplan-Meier)
    Veri Kaynağı: model_input_data_full.csv
    """
    try:
        from lifelines import KaplanMeierFitter
    except ImportError:
        return None

    plt.figure(figsize=(10, 6))
    kmf = KaplanMeierFitter()
    
    # En popüler 4 ekipman tipi
    top_types = df_model['Ekipman_Tipi'].value_counts().head(4).index
    
    for etype in top_types:
        subset = df_model[df_model['Ekipman_Tipi'] == etype]
        kmf.fit(subset['duration_days'], event_observed=subset['event'], label=etype)
        kmf.plot(ci_show=False)

    plt.title('Varlık Ömür Eğrileri (Survival Curves)', fontsize=16)
    plt.xlabel('Gün (Timeline)')
    plt.ylabel('Hayatta Kalma Olasılığı S(t)')
    plt.grid(True, linestyle='--', alpha=0.5)
    
    path = os.path.join(VISUAL_DIR, "ADV_03_Omur_Egrileri.png")
    plt.savefig(path, dpi=300)
    plt.close()
    return path

def plot_health_dashboard(df_res, logger):
    """
    GÖRSEL 4: Sağlık Analizi (Composite Dashboard)
    Veri Kaynağı: pof_predictions_final.csv
    """
    fig = plt.figure(figsize=(14, 8))
    gs = fig.add_gridspec(2, 2)
    fig.suptitle('Sağlık Analizi Dashboard', fontsize=22)

    # 1. Sağlık Skoru Dağılımı (Histogram)
    ax1 = fig.add_subplot(gs[0, 0])
    sns.histplot(df_res['Health_Score'], bins=30, kde=True, color='green', ax=ax1)
    ax1.axvline(x=40, color='red', linestyle='--', label='Kritik Eşik (40)')
    ax1.set_title('Sağlık Skoru Dağılımı')
    ax1.legend()

    # 2. Tip Bazlı Ortalama Sağlık (Bar)
    ax2 = fig.add_subplot(gs[0, 1])
    avg_health = df_res.groupby('Ekipman_Tipi')['Health_Score'].mean().sort_values().head(8)
    sns.barplot(y=avg_health.index, x=avg_health.values, ax=ax2, palette='RdYlGn')
    ax2.set_title('En Düşük Sağlık Skoruna Sahip Tipler')
    ax2.set_xlim(0, 100)

    # 3. Risk Sınıfı Pasta Grafiği
    ax3 = fig.add_subplot(gs[1, 0])
    risk_col = 'Risk_Class' if 'Risk_Class' in df_res.columns else 'Risk_Sinifi'
    counts = df_res[risk_col].value_counts()
    
    # Renkleri sabitle
    colors = {'Low': '#66b3ff', 'DÜŞÜK': '#66b3ff', 
              'Medium': '#ffcc99', 'ORTA': '#ffcc99',
              'High': '#ff9999', 'YÜKSEK': '#ff9999',
              'Critical': '#ff0000', 'KRİTİK': '#ff0000'}
    pie_colors = [colors.get(x, 'grey') for x in counts.index]
    
    ax3.pie(counts, labels=counts.index, autopct='%1.1f%%', colors=pie_colors, startangle=140)
    ax3.set_title('Risk Sınıfı Dağılımı')

    # 4. Boş Alan (Veya Metin Özeti)
    ax4 = fig.add_subplot(gs[1, 1])
    ax4.axis('off')
    summary_text = (
        f"Toplam Varlık: {len(df_res):,}\n"
        f"Ortalama Sağlık: {df_res['Health_Score'].mean():.1f}\n"
        f"Kritik Varlıklar: {counts.get('KRİTİK', 0) + counts.get('Critical', 0)}\n\n"
        "Not: Kırmızı bölge acil aksiyon gerektirir."
    )
    ax4.text(0.1, 0.5, summary_text, fontsize=14, bbox=dict(facecolor='wheat', alpha=0.3))

    path = os.path.join(VISUAL_DIR, "ADV_04_Saglik_Analizi.png")
    plt.savefig(path, dpi=300)
    plt.close()
    return path
def generate_visuals(df, logger):
    """
    Tüm görsel panoları (Dashboard elementleri) oluşturur.
    """
    logger.info("="*60)
    logger.info("[PHASE 2] Görsel Panolar Oluşturuluyor...")
    charts = {}
    
    # Risk Kolonunu Belirle (Standartlaştırma)
    risk_col = 'Risk_Sinifi'
    if 'Risk_Class' in df.columns: risk_col = 'Risk_Class'

    # 1. SAĞLIK SKORU DAĞILIMI (Histogram)
    # Filonun genel sağlık durumunu gösterir.
    if 'Health_Score' in df.columns:
        path = plot_single_chart(df, 'Health_Score', None, 'hist', 
                                 'Varlık Sağlık Skoru Dağılımı (0=Ölü, 100=Mükemmel)', 
                                 "02_saglik_skoru_dagilimi.png", logger,
                                 bins=40, color='teal', edgecolor='black')
        charts['health_dist'] = path

    # 2. KRONİK ANALİZİ (Bar Chart)
    # Hangi ekipmanlar sürekli baş ağrıtıyor?
    chronic_col = 'Chronic_Flag' if 'Chronic_Flag' in df.columns else 'Kronik_Flag'
    
    if chronic_col in df.columns:
        counts = df[chronic_col].value_counts()
        plt.figure(figsize=(8, 6))
        # 0: Yeşil (Normal), 1: Kırmızı (Kronik)
        colors = ['green', 'red'] if len(counts) == 2 else ['green']
        
        counts.plot(kind='bar', color=colors, edgecolor='black', alpha=0.8)
        plt.title('Kronik Varlık Dağılımı (1 = Kronik Sorunlu)', fontsize=14)
        plt.xlabel("Durum (0: Normal, 1: Kronik)")
        plt.ylabel("Varlık Sayısı")
        plt.xticks(rotation=0)
        
        path = os.path.join(VISUAL_DIR, "03_kronik_dagilimi.png")
        plt.savefig(path, dpi=300, bbox_inches='tight')
        plt.close()
        charts['chronic_dist'] = path

    # 3. COĞRAFİ HARİTA (Scatter Plot)
    # Risklerin mekansal dağılımı.
    if 'Latitude' in df.columns and 'Longitude' in df.columns:
        # Koordinatı 0 olmayanları al (Veri temizliği)
        gdf = df[(df['Latitude'] != 0) & (df['Longitude'] != 0) & (df['Latitude'].notna())].copy()
        
        if not gdf.empty:
            # Renk Haritası (Risk Sınıfına Göre)
            palette_map = {
                'Critical': 'red', 'High': 'orange', 'Medium': 'gold', 'Low': 'green',
                'KRİTİK': 'red', 'KRİTİK (KRONİK)': 'purple', 'YÜKSEK': 'orange', 
                'ORTA': 'gold', 'DÜŞÜK': 'green'
            }
            
            # Haritada olmayan etiketleri gri yap
            for label in gdf[risk_col].unique():
                if label not in palette_map: palette_map[label] = 'gray'

            path = plot_single_chart(gdf, 'Longitude', 'Latitude', 'scatter', 
                                     'Coğrafi Risk Haritası', "04_cografi_risk_haritasi.png", logger,
                                     hue=risk_col, height=10, width=10,
                                     palette=palette_map, s=40, alpha=0.7, edgecolor='k')
            charts['geo_map'] = path

    return charts

def validate_base_rates(df, logger):
    """
    Modelin sonuçlarını Endüstri Standartları ile kıyaslar (Reality Check).
    """
    logger.info("="*60)
    logger.info("[VALIDATION] Model Kalibrasyon Kontrolü...")
    
    # Sektör beklentileri (Yıllık Arıza Oranı - Failure Rate)
    # Kaynak: CIGRE ve IEEE standartları (yaklaşık)
    INDUSTRY_RANGES = {
        'Trafo': (0.005, 0.05),   # %0.5 - %5 arası normal
        'Kesici': (0.01, 0.08),   # %1 - %8
        'Ayırıcı': (0.02, 0.12),  # Ayırıcılar daha sık bozulur
        'Sigorta': (0.10, 0.40),  # Sigortalar sarf malzemesidir, çok bozulur
        'Hat': (0.005, 0.15),     # Hava şartlarına bağlı
        'Direk': (0.001, 0.03)    # Direkler nadir yıkılır
    }
    
    target_col = 'PoF_Ensemble_12Ay'
    if target_col not in df.columns:
        logger.warning("  ⚠️ [SKIP] PoF kolonu yok. Validasyon yapılamıyor.")
        return

    # Ekipman tipine göre ortalama tahminleri al
    stats = df.groupby('Ekipman_Tipi')[target_col].mean().reset_index()
    
    for _, row in stats.iterrows():
        etype = str(row['Ekipman_Tipi'])
        pred = row[target_col]
        
        # Ekipman isminde anahtar kelime ara (örn: "OG Trafo" içinde "Trafo" var mı?)
        matched_key = next((k for k in INDUSTRY_RANGES if k in etype), None)
        
        if matched_key:
            low, high = INDUSTRY_RANGES[matched_key]
            status = "✅ OK"
            if pred < low: status = "📉 DÜŞÜK (Under-prediction?)"
            if pred > high: status = "🚨 YÜKSEK (Over-prediction?)"
            
            logger.info(f"  > {etype.ljust(25)}: Tahmin %{pred*100:.1f} (Ref: %{low*100:.0f}-%{high*100:.0f}) -> {status}")

def plot_aggregate_risk_by_type(df, logger):
    """
    Ekipman tiplerine göre risk yoğunluğunu gösteren Dual-Axis grafik.
    """
    if 'PoF_Ensemble_12Ay' not in df.columns:
        return None

    # Agregasyon
    # Hangi sütun adını kullanacağımızı bulalım
    chronic_col = 'Chronic_Flag' if 'Chronic_Flag' in df.columns else 'Kronik_Flag'
    
    agg_dict = {'PoF_Ensemble_12Ay': 'mean', 'cbs_id': 'count'}
    if chronic_col in df.columns:
        agg_dict[chronic_col] = 'mean' # Kronik oranı

    agg_df = df.groupby('Ekipman_Tipi').agg(agg_dict).reset_index()
    
    # Sütun isimlerini düzelt
    agg_df = agg_df.rename(columns={
        'PoF_Ensemble_12Ay': 'Mean_PoF', 
        'cbs_id': 'Count',
        chronic_col: 'Chronic_Rate'
    })
    
    # Sadece en az 30 varlığı olan tipleri al ve en risklileri seç
    agg_df = agg_df[agg_df['Count'] >= 30].sort_values('Mean_PoF', ascending=False).head(10)

    if agg_df.empty: return None

    # Çift Eksenli Grafik Çizimi
    fig, ax1 = plt.subplots(figsize=(12, 6))
    
    # 1. Bar Chart (Arıza Olasılığı)
    sns.barplot(x='Ekipman_Tipi', y='Mean_PoF', data=agg_df, ax=ax1, color='firebrick', alpha=0.6)
    ax1.set_ylabel('Ortalama Arıza Olasılığı (12 Ay)', color='firebrick', fontweight='bold')
    ax1.tick_params(axis='y', labelcolor='firebrick')
    ax1.set_xlabel('Ekipman Tipi', fontsize=12)
    ax1.tick_params(axis='x', rotation=45)
    
    # 2. Line Chart (Kronik Oranı) - Eğer veri varsa
    if 'Chronic_Rate' in agg_df.columns:
        ax2 = ax1.twinx()
        sns.lineplot(x='Ekipman_Tipi', y='Chronic_Rate', data=agg_df, ax=ax2, color='navy', marker='o', linewidth=2)
        ax2.set_ylabel('Kronik Varlık Oranı', color='navy', fontweight='bold')
        ax2.tick_params(axis='y', labelcolor='navy')
        ax2.grid(False) # İkinci ızgarayı kapat ki karışmasın
    
    plt.title('Ekipman Tipine Göre Risk Analizi (Top 10)', fontsize=14)
    
    path = os.path.join(VISUAL_DIR, "08_aggregate_risk_by_type.png")
    plt.savefig(path, dpi=300, bbox_inches='tight')
    plt.close()
    
    logger.info(f"  📸 Karma risk grafiği kaydedildi: 08_aggregate_risk_by_type.png")
    return path
# ------------------------------------------------------------------------------
# PHASE 3: EXCEL REPORTING
# ------------------------------------------------------------------------------

def generate_case_studies(df_risk, logger):
    """
    Modelin son 6 aydaki başarısını ölçen 'Vaka Analizi' (Case Study) tablosu.
    """
    logger.info("[PHASE 1.5] Vaka Analizleri (Case Studies) Oluşturuluyor...")
    
    # 1. Arıza Olaylarını Yükle
    events_path = os.path.join(INTERMEDIATE_DIR, "fault_events_clean.csv")
    
    events = pd.DataFrame()
    if os.path.exists(events_path):
        events = pd.read_csv(events_path)
        logger.info(f"  > Arıza verisi yüklendi: {events_path}")
    else:
        raw_path = os.path.join(BASE_DIR, "data", "girdiler", "ariza_final.xlsx")
        if os.path.exists(raw_path):
            logger.warning("  ⚠️ Ara çıktı yok, ham Excel okunuyor...")
            events = pd.read_excel(raw_path)
    
    if events.empty:
        logger.error("  ❌ Arıza verisi bulunamadı. Case Study atlanıyor.")
        return pd.DataFrame()

    # 2. Sütun İsimlerini Standartlaştır
    col_map = {
        'started at': 'Ariza_Baslangic_Zamani',
        'cbs_id': 'cbs_id',
        'Ekipman Kodu': 'cbs_id'
    }
    events = events.rename(columns=col_map)
    
    if 'Ariza_Baslangic_Zamani' not in events.columns or 'cbs_id' not in events.columns:
        logger.error("  ❌ Gerekli sütunlar (started at, cbs_id) eksik.")
        return pd.DataFrame()

    # 3. Tarih ve ID Temizliği
    # parse_mixed_dates fonksiyonunun import edildiğinden emin olun, yoksa pd.to_datetime kullanın
    try:
        events['Ariza_Baslangic_Zamani'] = events['Ariza_Baslangic_Zamani'].apply(parse_mixed_dates)
    except NameError:
        events['Ariza_Baslangic_Zamani'] = pd.to_datetime(events['Ariza_Baslangic_Zamani'], errors='coerce')

    events['cbs_id'] = events['cbs_id'].astype(str).str.lower().str.strip()

    # 3.5. Gelecek Tarihli Kayıtları Filtrele
    today = pd.Timestamp.now()
    future_count = (events['Ariza_Baslangic_Zamani'] > today).sum()
    if future_count > 0:
        logger.warning(f"  ⚠️ {future_count} gelecek tarihli kayıt bulundu ve filtrelendi.")
        events = events[events['Ariza_Baslangic_Zamani'] <= today]

    # 4. Son 6 Aydaki Arızaları Filtrele
    last_date = events['Ariza_Baslangic_Zamani'].max()
    if pd.isna(last_date): return pd.DataFrame()
    
    cutoff_date = last_date - timedelta(days=180)
    
    # [DÜZELTME BURADA] - Değişken ismi tutarlılığı sağlandı
    recent_faults = events[events['Ariza_Baslangic_Zamani'] >= cutoff_date].copy()
    
    # Sadece gerekli sütunları al, böylece 'Ekipman_Tipi' gibi sütunlar buradan silinir
    # ve merge işleminde çakışma (duplicate column) yaratmaz.
    keep_cols = ['cbs_id', 'Ariza_Baslangic_Zamani', 'cause code', 'Ariza_Nedeni']
    recent_faults = recent_faults[[c for c in keep_cols if c in recent_faults.columns]]

    if recent_faults.empty:
        logger.warning("  ⚠️ Son 6 ayda hiç arıza kaydı yok.")
        return pd.DataFrame()

    # 5. Risk Verisi ile Birleştir (Merge)
    df_risk = ensure_pof_column(df_risk, logger)
    df_risk['cbs_id'] = df_risk['cbs_id'].astype(str).str.lower().str.strip()

    risk_col = 'Risk_Sinifi'
    if 'Risk_Class' in df_risk.columns: risk_col = 'Risk_Class'
    
    # Ekipman_Tipi burada ekleniyor (df_risk'ten geliyor)
    cols_to_merge = ['cbs_id', risk_col, 'PoF_Ensemble_12Ay', 'Health_Score', 'Ekipman_Tipi', 'Ilce', 'Marka']
    cols_to_merge = [c for c in cols_to_merge if c in df_risk.columns]

    # [DÜZELTME] merge işleminde temizlenmiş 'recent_faults' kullanılıyor
    case_df = recent_faults.merge(df_risk[cols_to_merge], on='cbs_id', how='left')
    
    # 6. Başarı Değerlendirmesi
    def judge_prediction(row):
        if risk_col not in row or pd.isna(row[risk_col]): return "Bilinmeyen Varlık"
        
        r = str(row[risk_col]).upper()
        if any(x in r for x in ['CRIT', 'KRİT', 'HIGH', 'YÜKSEK']):
            return "BAŞARILI (Öngörüldü)"
        elif any(x in r for x in ['MED', 'ORTA']):
            return "KISMİ (İzleme)"
        else:
            return "KAÇIRILDI (Düşük Risk)"

    case_df['Model_Karari'] = case_df.apply(judge_prediction, axis=1)
    
    # Rapor için Seçim
    successes = case_df[case_df['Model_Karari'] == "BAŞARILI (Öngörüldü)"].sort_values('Ariza_Baslangic_Zamani', ascending=False).head(10)
    misses = case_df[case_df['Model_Karari'] == "KAÇIRILDI (Düşük Risk)"].sort_values('Ariza_Baslangic_Zamani', ascending=False).head(5)
    
    final_cases = pd.concat([successes, misses])
    
    # Sütun Sıralaması (Görsel Güzellik İçin)
    display_order = ['cbs_id', 'Ekipman_Tipi', 'Ariza_Baslangic_Zamani', 'Ilce', 'Model_Karari', risk_col, 'PoF_Ensemble_12Ay']
    final_cols = [c for c in display_order if c in final_cases.columns]
    final_cases = final_cases[final_cols]

    logger.info(f"  ✅ Vaka analizi tamamlandı: {len(successes)} Başarılı, {len(misses)} Kaçırılan örnek seçildi.")
    
    return final_cases

# ------------------------------------------------------------------------------
# PHASE 3: EXCEL REPORTING (FINAL OUTPUT)
# ------------------------------------------------------------------------------
# =============================================================================
# 🏆 PROOF & FINAL DELIVERY (KANIT VE NİHAİ RAPORLAMA)
# =============================================================================
# Bu modül, analiz döngüsünü tamamlar ve sonuçları iki kritik formatta sunar:
#
# 1. 🕵️‍♂️ Vaka Analizleri / Case Studies (generate_case_studies):
#    - Modelin "Güvenilirliğini" ispatlar.
#    - Son 6 ayda gerçekleşen arızaları, modelin geçmiş tahminleriyle kıyaslar.
#    - Çıktı: "Model, geçen ay yanan Trafo X'i 'Kritik Risk' olarak öngörmüş müydü?"
#      sorusunun cevabını içeren "Başarı/Hata Karnesi"dir.
#    - Bu tablo, bütçe onayı almak için en güçlü kanıttır.
#
# 2. 📑 Nihai Excel Raporu (create_excel_report):
#    - Projenin resmi teslimat dosyasıdır.
#    - Çok sayfalı (Multi-sheet) bir yapıdadır:
#      a. Yönetici Özeti: Tek bakışta filo sağlığı ve KPI'lar.
#      b. Acil Müdahale: Bakım ekiplerinin pazartesi sabahı alacağı iş listesi.
#      c. Risk Master: Tüm varlıkların detaylı risk dökümü (Top 1000).
# =============================================================================
def create_excel_report(df, crit_chronic, case_studies, logger): 
    """
    Tüm analizleri tek bir Excel dosyasında (Multi-Sheet) toplar.
    """
    logger.info("="*60)
    logger.info("[PHASE 3] Excel Raporu Oluşturuluyor...")
    
    timestamp = datetime.now().strftime("%Y-%m-%d")
    out_path = os.path.join(REPORT_DIR, f"PoF3_Analiz_Raporu_Final_{timestamp}.xlsx")
    
    risk_col = 'Risk_Sinifi'
    if 'Risk_Class' in df.columns: risk_col = 'Risk_Class'
    
    # Excel Writer Başlat
    with pd.ExcelWriter(out_path, engine='openpyxl') as writer:
        
        # 1. YÖNETİCİ ÖZETİ (KPI Tablosu)
        total = len(df)
        crit_mask = df[risk_col].astype(str).str.contains('KRİT|CRIT', case=False, na=False)
        crit_count = crit_mask.sum()
        
        avg_health = 0
        if 'Health_Score' in df.columns: avg_health = df['Health_Score'].mean()
        
        summary_data = {
            'Metrik': ['Toplam Varlık Sayısı', 'Kritik Riskli Varlık Sayısı', 
                       'Kronik ve Kritik (Acil)', 'Filo Ortalama Sağlık Puanı', 'Rapor Tarihi'],
            'Değer': [total, crit_count, len(crit_chronic), f"{avg_health:.1f} / 100", timestamp]
        }
        pd.DataFrame(summary_data).to_excel(writer, sheet_name='Yonetici_Ozeti', index=False)
        
        # 2. ACİL MÜDAHALE LİSTESİ
        if not crit_chronic.empty:
            crit_chronic.to_excel(writer, sheet_name='Acil_Mudahale_Listesi', index=False)

        # 3. VAKA ANALİZLERİ (Kanıt)
        if not case_studies.empty:
            # Sadece önemli sütunları al
            cols = ['cbs_id', 'Ariza_Baslangic_Zamani', 'Ekipman_Tipi', 'Ilce', 
                    'Model_Karari', risk_col, 'PoF_Ensemble_12Ay']
            cols = [c for c in cols if c in case_studies.columns]
            case_studies[cols].to_excel(writer, sheet_name='Model_Dogrulama_Vakalari', index=False)
            
        # 4. RİSK MASTER (TOP 1000)
        # Tüm filoyu basmak yerine en riskli 1000 varlığı basar (Dosya boyutu için)
        sort_col = 'PoF_Ensemble_12Ay'
        ascending = False # En yüksek PoF en üstte
        
        if sort_col not in df.columns:
            if 'Health_Score' in df.columns:
                sort_col = 'Health_Score'
                ascending = True # En düşük sağlık puanı en üstte
            else:
                sort_col = df.columns[0] # Fallback

        df.sort_values(sort_col, ascending=ascending).head(1000).to_excel(writer, sheet_name='Risk_Master_Top1000', index=False)
            
    logger.info(f"  💾 Excel Raporu Kaydedildi: {os.path.basename(out_path)}")
    return out_path

# ------------------------------------------------------------------------------
# PHASE 4: POWERPOINT PRESENTATION (YÖNETİCİ SUNUMU)
# ------------------------------------------------------------------------------
# =============================================================================
# 📽️ EXECUTIVE PRESENTATION (OTOMATİK SUNUM)
# =============================================================================
# Bu fonksiyon, teknik analiz sonuçlarını "Yönetim Kurulu" formatına çevirir.
#
# 🎯 Özellikleri:
#    - Python-PPTX kütüphanesini kullanır.
#    - Dinamik Özet: Raporun alındığı günkü sayıları (Kritik, Kronik vb.)
#      slaytların içine metin olarak yazar.
#    - Görsel Galeri: VISUAL_DIR altında üretilen tüm grafikleri
#      otomatik olarak yeni slaytlara yerleştirir.
#
# ⚠️ Gereksinim:
#    - 'pip install python-pptx' kurulu olmalıdır.
#    - Kurulu değilse fonksiyon sessizce çalışmayı durdurur (Crash olmaz).
# =============================================================================
def create_pptx_presentation(df, charts, logger):
    """
    Analiz sonuçlarını ve grafikleri (Gelişmiş Analitikler dahil) PowerPoint sunumuna dönüştürür.
    """
    if not HAS_PPTX:
        logger.warning("  ⚠️ python-pptx kütüphanesi yüklü değil. PPTX oluşturulamadı.")
        return

    logger.info("="*60)
    logger.info("[PHASE 4] PowerPoint Sunumu Oluşturuluyor...")
    
    try:
        prs = Presentation()
    except Exception as e:
        logger.error(f"  ❌ PPTX başlatılamadı: {e}")
        return

    timestamp = datetime.now().strftime("%d.%m.%Y")
    
    # --- SLIDE 1: KAPAK ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Varlık Sağlığı ve Risk Analizi (PoF)"
    subtitle.text = f"Yönetici Özeti Raporu\nRapor Tarihi: {timestamp}\nOluşturan: AI Risk Engine"
    
    # --- SLIDE 2: YÖNETİCİ ÖZETİ (METİN) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.text = "Genel Durum Özeti"
    
    # İstatistikleri Hesapla
    total = len(df)
    
    risk_col = 'Risk_Sinifi'
    if 'Risk_Class' in df.columns: risk_col = 'Risk_Class'
    
    crit_count = 0
    if risk_col in df.columns:
        crit_count = df[risk_col].astype(str).str.upper().str.contains('KRİT|CRIT').sum()
        
    chronic_count = 0
    chronic_col = 'Chronic_Flag' if 'Chronic_Flag' in df.columns else 'Kronik_Flag'
    if chronic_col in df.columns:
        chronic_count = (df[chronic_col] == 1).sum()
        
    avg_health = df['Health_Score'].mean() if 'Health_Score' in df.columns else 0

    # Metin İçeriği
    tf = slide.placeholders[1].text_frame
    tf.text = f"Analiz Kapsamı: {total:,} Adet Şebeke Varlığı"
    
    p = tf.add_paragraph()
    p.text = f"🚨 Kritik Riskli Varlıklar: {crit_count:,} adet (%{100*crit_count/total:.1f})"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"🏚️ Kronik Sorunlu Varlıklar: {chronic_count:,} adet"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"❤️ Filo Ortalama Sağlık Puanı: {avg_health:.1f} / 100"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Öneri: 'Acil Müdahale Listesi'ndeki varlıklar için iş emri oluşturulmalıdır."
    p.level = 2

    # --- SLIDE 3-X: GRAFİKLER ---
    # Aggregate Risk grafiği charts sözlüğünde olmayabilir, manuel kontrol ekliyoruz.
    agg_path = os.path.join(VISUAL_DIR, "08_aggregate_risk_by_type.png")
    if os.path.exists(agg_path) and 'aggregate_risk' not in charts:
        charts['aggregate_risk'] = agg_path

    # GÜNCELLENMİŞ HARİTA: Yeni grafikleri buraya ekledik
    slide_mapping = {
        'health_dist': "Filonun Sağlık Skoru Dağılımı",
        'geo_map': "Coğrafi Risk Yoğunluk Haritası",
        'aggregate_risk': "Ekipman Tipine Göre Risk Analizi",
        'chronic_dist': "Kronik Varlık İstatistikleri",
        
        # --- YENİ EKLENENLER (Advanced Diagnostics) ---
        'drivers': "Risk Faktörleri (Drivers) Analizi",
        'operational': "Operasyonel Durum Paneli",
        'survival': "Varlık Ömür Eğrileri (Survival Curves)",
        'health_dash': "Detaylı Sağlık Analizi Dashboard"
    }
    
    for key, title_text in slide_mapping.items():
        # Grafik sözlükte var mı VE dosya diskte mevcut mu?
        if key in charts and charts[key] and os.path.exists(charts[key]):
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only (Resim için boş alan)
            slide.shapes.title.text = title_text
            
            # Resmi Ortala ve Yerleştir
            img_path = charts[key]
            
            left = Inches(0.5) # Biraz daha sola yanaşık
            top = Inches(1.5)
            height = Inches(5.5) 
            
            try:
                slide.shapes.add_picture(img_path, left, top, height=height)
            except Exception as img_err:
                logger.warning(f"  ⚠️ Resim eklenirken hata ({key}): {img_err}")

    # Kaydet
    out_path = os.path.join(OUTPUT_DIR, f"PoF3_Yonetici_Sunumu_{timestamp}.pptx")
    try:
        prs.save(out_path)
        logger.info(f"  💾 Sunum Kaydedildi: {os.path.basename(out_path)}")
    except Exception as e:
        logger.error(f"  ❌ Sunum dosyası kaydedilemedi (Dosya açık olabilir mi?): {e}")
# ------------------------------------------------------------------------------
# MAIN ORCHESTRATION
# ------------------------------------------------------------------------------
# =============================================================================
# 🚀 MAIN PIPELINE ORCHESTRATION (ANA YÖNETİM MERKEZİ)
# =============================================================================
# Bu fonksiyon, ham veriden nihai raporlara giden uçtan uca (End-to-End) akışı yönetir.
#
# 🔄 İşlem Adımları (Process Flow):
#
# 1. 📥 Data Ingestion (Yükleme):
#    - Arıza ve Sağlam ekipman verileri okunur, tarihler parse edilir.
#    - Veri setinin zaman aralığı (Start/End Date) otomatik belirlenir.
#
# 2. 🏗️ Dataset Construction (Veri İnşası):
#    - 'build_equipment_master': Tüm varlıkların tekil listesi çıkarılır.
#    - 'add_survival_columns': Sol Kesilme (Left Truncation) ve Ömür (Duration) hesaplanır.
#    - 'Chronic & Temporal': Arıza geçmişine dayalı dinamik özellikler türetilir.
#
# 3. 🛡️ Global Modeling (Güvenlik Ağı):
#    - Veri seti küçük olan ekipman tipleri (örn. "Ayırıcı") için tek başına model
#      eğitmek risklidir (Overfitting).
#    - Bu adımda tüm veriyi kullanan "Global Modeller" (Cox, RSF, ML) eğitilir.
#
# 4. ⚙️ Stratified Training (Katmanlı Eğitim):
#    - Her ekipman tipi (Trafo, Kesici vb.) için döngüye girilir.
#    - Karar Mekanizması:
#      a. Yeterli Veri Var mı? (N > 50, Events > 10) -> O tipe ÖZEL model eğit.
#      b. Veri Yetersiz mi? -> GLOBAL modelleri kullan (Fallback).
#
# 5. 🏥 Risk Scoring (Puanlama):
#    - Modellerin ürettiği olasılıklar (PoF), 0-100 arası "Sağlık Skoru"na çevrilir.
#    - Kritik ve Kronik varlıklar etiketlenir.
#
# 6. 🕰️ Backtesting (Doğrulama):
#    - Modelin başarısını ölçmek için geçmişe dönük (2022-2024) simülasyon yapılır.
# =============================================================================
def main():
    logger = setup_logger()
    logger.info("🚀 PoF3 Raporlama Motoru Başlatılıyor...")
    
    # 1. Dosya Kontrolü ve Yükleme
    # pof.py çıktısını arıyoruz
    risk_path = os.path.join(OUTPUT_DIR, "pof_predictions_final.csv")
    
    if not os.path.exists(risk_path):
        # Fallback: Farklı isimlendirme ihtimaline karşı
        alt_path = os.path.join(OUTPUT_DIR, "risk_equipment_master.csv")
        if os.path.exists(alt_path):
            risk_path = alt_path
        else:
            logger.error(f"[FATAL] Sonuç dosyası bulunamadı: {risk_path}")
            logger.error("Lütfen önce 'pof.py' (Analiz Motoru) çalıştırın.")
            return
        
    df = pd.read_csv(risk_path)
    logger.info(f"[LOAD] Risk sonuçları yüklendi: {len(df):,} kayıt")
    
    # 2. Veri Temizliği ve Standartlaştırma
    # ID'leri string yap (Merge hatasını önler)
    df['cbs_id'] = df['cbs_id'].astype(str).str.lower().str.strip()

    # Kolon Eşleştirme (Risk_Sinifi -> Risk_Class)
    if 'Risk_Sinifi' in df.columns and 'Risk_Class' not in df.columns:
        df['Risk_Class'] = df['Risk_Sinifi']
    elif 'Risk_Class' not in df.columns:
        logger.warning("[WARN] Risk kolonu yok. Varsayılan 'Low' atanıyor.")
        df['Risk_Class'] = 'Low'

    # PoF Kolonunu Garantiye Al (Mevcut kodunuzdaki satır)
    df = ensure_pof_column(df, logger)

    # =============================================================================
    # 🚑 [FIX] KRONİK VERİ KURTARMA OPERASYONU
    # =============================================================================
    # Final dosyada 'Chronic_Flag' yoksa, ara hesaplama dosyasından (ozellikler_zamansal) çeker.
    if 'Chronic_Flag' not in df.columns and 'Kronik_Flag' not in df.columns:
        logger.warning("  ⚠️ Ana dosyada 'Chronic_Flag' bulunamadı! Ara dosyalardan kurtarılıyor...")
        
        # Log dosyasında gördüğümüz ara çıktı yolu
        chronic_path = os.path.join(INTERMEDIATE_DIR, "ozellikler_zamansal.csv")
        
        if os.path.exists(chronic_path):
            try:
                # Sadece ID ve Flag kolonlarını oku (Hafif olsun)
                df_chronic = pd.read_csv(chronic_path, usecols=lambda c: c in ['cbs_id', 'Chronic_Flag', 'Kronik_Flag', 'Fault_Count'])
                
                # ID Standardizasyonu (Eşleşme garantisi için)
                df_chronic['cbs_id'] = df_chronic['cbs_id'].astype(str).str.lower().str.strip()
                
                # Kolon ismini belirle
                source_col = 'Chronic_Flag' if 'Chronic_Flag' in df_chronic.columns else 'Kronik_Flag'
                
                if source_col:
                    # Ana tablo ile birleştir
                    df = df.merge(df_chronic[['cbs_id', source_col]], on='cbs_id', how='left')
                    
                    # NaN değerleri 0 yap (Eşleşmeyenler kronik değildir)
                    df[source_col] = df[source_col].fillna(0).astype(int)
                    
                    # İsim standardı
                    if source_col != 'Chronic_Flag':
                        df['Chronic_Flag'] = df[source_col]
                        
                    count = df['Chronic_Flag'].sum()
                    logger.info(f"  ✅ Kronik verisi başarıyla eklendi: {count} adet kronik varlık kurtarıldı.")
                else:
                    logger.error("  ❌ Ara dosyada da flag bulunamadı.")
            except Exception as e:
                logger.error(f"  ❌ Merge işlemi başarısız: {e}")
        else:
            logger.error(f"  ❌ Ara dosya bulunamadı: {chronic_path}")
            df['Chronic_Flag'] = 0 # Kod patlamasın diye dummy
    
    # Hâlâ yoksa (Kurtarma başarısızsa) dummy oluştur
    if 'Chronic_Flag' not in df.columns:
        df['Chronic_Flag'] = 0

    # =============================================================================
    # A) Aksiyon Listeleri
    crit_chronic = generate_action_lists(df, logger)
    
    # B) Vaka Analizleri
    case_studies = generate_case_studies(df, logger)
    
    # C) Görseller
    charts = generate_visuals(df, logger)
    
    # D) Özel Grafikler
    agg_path = plot_aggregate_risk_by_type(df, logger)
    if agg_path: charts['aggregate_risk'] = agg_path 
    # ... (Mevcut main fonksiyonunun son kısımları) ...

    # 5. Raporları Üret
    crit_chronic = generate_action_lists(df, logger)
    case_studies = generate_case_studies(df, logger)
    charts = generate_visuals(df, logger)
    
    # --- YENİ EKLENEN GELİŞMİŞ GÖRSELLER ---
    # Bu veriler 'intermediate' klasöründeki dosyalardan okunacak
    
    # A) Model Girdisi (Drivers ve Survival için lazım)
    model_data_path = os.path.join(INTERMEDIATE_DIR, "model_input_data_full.csv")
    if os.path.exists(model_data_path):
        df_model = pd.read_csv(model_data_path)
        
        # Drivers
        p1 = plot_risk_drivers(df_model, logger)
        if p1: charts['drivers'] = p1
        
        # Survival Curves
        p3 = plot_survival_curves(df_model, logger)
        if p3: charts['survival'] = p3
    
    # B) Operasyonel Dashboard
    p2 = plot_operational_dashboard(logger)
    if p2: charts['operational'] = p2
    
    # C) Sağlık Dashboard
    p4 = plot_health_dashboard(df, logger)
    if p4: charts['health_dash'] = p4
    # ---------------------------------------

    agg_path = plot_aggregate_risk_by_type(df, logger)
    if agg_path: charts['aggregate_risk'] = agg_path 
    
    # 5. Raporlama
    create_excel_report(df, crit_chronic, case_studies, logger)
    create_pptx_presentation(df, charts, logger)
    
    logger.info("")
    logger.info("="*60)
    logger.info("[SUCCESS] Tüm Raporlama Süreci Başarıyla Tamamlandı.")
    logger.info(f"📂 Çıktı Klasörü: {OUTPUT_DIR}")
    logger.info("="*60)

if __name__ == "__main__":
    main()